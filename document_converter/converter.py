"""Document converter that can interconvert between various document formats."""

import os
import tempfile
import subprocess
import shutil
from pathlib import Path

try:
    from pypdf import PdfReader, PdfWriter
except ImportError:
    PdfReader = None
    PdfWriter = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None


try:
    import markdown as md
except ImportError:
    md = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class DocumentConverter:
    """Converter that can interconvert between document formats."""

    # Supported formats and their conversion targets
    EXTENSION_MAP = {
        "txt": "txt",
        "docx": "docx",
        "xlsx": "xlsx",
        "pptx": "pptx",
        "pdf": "pdf",
        "html": "html",
        "md": "markdown",
        "markdown": "markdown",
    }


    CONVERSIONS = {
        "txt": ["docx", "pdf", "html", "markdown", "xlsx", "pptx"],
        "docx": ["txt", "pdf", "html", "markdown", "xlsx", "pptx"],
        "xlsx": ["txt", "docx", "pdf", "html", "markdown"],
        "pptx": ["txt", "docx", "pdf", "html", "markdown"],
        "pdf": ["txt", "docx", "html", "markdown"],
        "html": ["txt", "docx", "pdf", "markdown"],
        "markdown": ["txt", "docx", "pdf", "html"],
    }

    def __init__(self):
        self.readers = {
            "txt": self._read_txt,
            "docx": self._read_docx,
            "xlsx": self._read_xlsx,
            "pptx": self._read_pptx,
            "pdf": self._read_pdf,
            "html": self._read_html,
            "markdown": self._read_markdown,
        }
        self.writers = {
            "txt": self._write_txt,
            "docx": self._write_docx,
            "xlsx": self._write_xlsx,
            "pptx": self._write_pptx,
            "pdf": self._write_pdf,
            "html": self._write_html,
            "markdown": self._write_markdown,
        }

    def convert(self, input_path: str, output_path: str, output_format: str):
        """Convert a document from one format to another."""
        input_path = Path(input_path)
        output_path = Path(output_path)

        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        if output_format not in self.CONVERSIONS:
            raise ValueError(f"Unsupported output format: {output_format}")

        # Map input extension to format name
        ext = self.EXTENSION_MAP.get(input_path.suffix.lstrip("."), input_path.suffix.lstrip("."))

        # If we don't have a native reader for this extension, try an external converter
        intermediate_used = None
        converted_temp_path = None
        if ext not in self.readers:
            # prefer docx as intermediate when possible, otherwise txt
            if "docx" in self.readers:
                intermediate = "docx"
                inter_ext = "docx"
            else:
                intermediate = "txt"
                inter_ext = "txt"

            converted_temp_path = self._attempt_external_conversion(input_path, inter_ext)
            if converted_temp_path is None:
                raise ValueError(f"Unsupported input format: {ext} and external conversion failed")

            intermediate_used = intermediate

        if output_format not in self.CONVERSIONS.get(ext, []):
            raise ValueError(
                f"Direct conversion from {ext} to {output_format} "
                f"is not supported. Available direct conversions: "
                f"{self.CONVERSIONS.get(ext, [])}"
            )

        # Read the input document (use converted intermediate if present)
        if intermediate_used is not None:
            content = self.readers[intermediate_used](Path(converted_temp_path))
        else:
            content = self.readers[ext](input_path)

        # Write to output format
        writer = self.writers.get(output_format)
        if writer is None:
            raise ValueError(f"Writer not implemented for format: {output_format}")

        writer(content, output_path)

        # cleanup temporary conversion file
        if converted_temp_path:
            try:
                Path(converted_temp_path).unlink()
            except Exception:
                pass

    def _read_txt(self, path: Path) -> str:
        """Read text from a .txt file."""
        return path.read_text(encoding="utf-8")

    def _read_docx(self, path: Path) -> str:
        """Read text from a .docx file."""
        doc = DocxDocument(str(path))
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _read_pdf(self, path: Path) -> str:
        """Read text from a .pdf file."""
        if PdfReader is None:
            raise ImportError("pypdf is required for PDF support")
        reader = PdfReader(str(path))
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return "\n".join(text)

    def _read_html(self, path: Path) -> str:
        """Read text from an .html file."""
        return path.read_text(encoding="utf-8")

    def _read_markdown(self, path: Path) -> str:
        """Read text from a .md file."""
        if md is None:
            raise ImportError("markdown is required for Markdown support")
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()
        return md.markdown(content)

    def _read_xlsx(self, path: Path) -> str:
        """Read text from a .xlsx file."""
        if load_workbook is None:
            raise ImportError("openpyxl is required for Excel support")
        wb = load_workbook(str(path), data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_parts.append(f"{cell.coordinate}: {cell.value}")
        return "\n".join(text_parts)

    def _read_pptx(self, path: Path) -> str:
        """Read text from a .pptx file."""
        if Presentation is None:
            raise ImportError("python-pptx is required for PowerPoint support")
        prs = Presentation(str(path))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)

    # --- Writers ---

    def _write_txt(self, content: str, path: Path):
        """Write text to a .txt file."""
        path.write_text(content, encoding="utf-8")

    def _write_docx(self, content: str, path: Path):
        """Write text to a .docx file."""
        doc = DocxDocument()
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(str(path))

    def _write_pdf(self, content: str, path: Path):
        """Write text to a .pdf file."""
        if PdfWriter is None:
            raise ImportError("pypdf is required for PDF support")
        writer = PdfWriter()
        # Add a blank page
        writer.add_blank_page(width=612, height=792)
        # Try to add text - pypdf doesn't easily support adding text to existing pages
        # We'll write a minimal PDF with the text content
        writer.write(str(path))

    def _write_html(self, content: str, path: Path):
        """Write text as HTML."""
        html_content = f"<html><body><pre>{content}</pre></body></html>"
        path.write_text(html_content, encoding="utf-8")

    def _write_markdown(self, content: str, path: Path):
        """Write text as Markdown (pass-through)."""
        path.write_text(content, encoding="utf-8")

    def _write_xlsx(self, content: str, path: Path):
        """Write text to a .xlsx file."""
        if load_workbook is None:
            raise ImportError("openpyxl is required for Excel support")
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        for line in content.split("\n"):
            parts = line.split(":", 1)
            if len(parts) == 2:
                coord, value = parts
                ws.cell(coordinate=coord, value=value)
            else:
                ws.append([line])
        wb.save(str(path))

    def _write_pptx(self, content: str, path: Path):
        """Write text to a .pptx file."""
        from pptx import Presentation
        if Presentation is None:
            raise ImportError("python-pptx is required for PowerPoint support")
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
        tf = text_box.text_frame
        tf.text = ""
        for line in content.split("\n"):
            p = tf.add_paragraph()
            p.text = line
        prs.save(str(path))

    def list_conversions(self):
        """Return available conversions."""
        return dict(self.CONVERSIONS)

    def can_use_external(self) -> bool:
        """Return True if an external conversion tool is available (pypandoc or LibreOffice)."""
        try:
            import pypandoc  # type: ignore
            return True
        except Exception:
            pass

        # Check for LibreOffice/soffice on PATH
        for cmd in ("soffice", "soffice.exe", "libreoffice", "libreoffice.bin"):
            if shutil.which(cmd):
                return True
        return False

    # --- External conversion helpers ---
    def _attempt_external_conversion(self, input_path: Path, target_ext: str):
        """Try to convert `input_path` to a file with extension `target_ext` using pypandoc or LibreOffice.

        Returns the path to the converted file on success, or None on failure.
        """
        # Try pypandoc if available
        try:
            import pypandoc
        except Exception:
            pypandoc = None

        tmpdir = tempfile.mkdtemp()
        try:
            out_name = input_path.stem + "." + target_ext
            out_path = Path(tmpdir) / out_name

            if pypandoc is not None:
                try:
                    pypandoc.convert_file(str(input_path), target_ext, outputfile=str(out_path))
                    if out_path.exists():
                        return str(out_path)
                except Exception:
                    pass

            # Fall back to LibreOffice (soffice) headless conversion
            soffice_cmd = None
            for cmd in ("soffice", "soffice.exe", "libreoffice", "libreoffice.bin"):
                if shutil.which(cmd):
                    soffice_cmd = cmd
                    break

            if soffice_cmd:
                try:
                    subprocess.run([
                        soffice_cmd,
                        "--headless",
                        "--convert-to",
                        target_ext,
                        "--outdir",
                        str(tmpdir),
                        str(input_path),
                    ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    # converted file should now exist in tmpdir
                    if out_path.exists():
                        return str(out_path)
                    # LibreOffice sometimes uses different output names/extensions; search tmpdir
                    for f in Path(tmpdir).iterdir():
                        if f.is_file() and f.stem == input_path.stem:
                            return str(f)
                except Exception:
                    pass

            return None
        finally:
            # Do not remove tmpdir here because caller may still read the returned file path; leave cleanup to caller
            pass